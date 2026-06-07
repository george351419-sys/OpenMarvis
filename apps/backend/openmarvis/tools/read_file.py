"""read_file: 复杂格式文档解析为纯文本（Markdown）。

read_text 只处理纯文本文件；read_file 按扩展名分派到专用解析器，
返回 LLM 友好的 Markdown 文本，支持 offset/limit 分页。

支持格式：.pdf .docx .pptx .xlsx .xlsm .csv .md .txt .json .yaml .yml
不支持的扩展名直接拒绝，让 LLM 自己处理（比如用 analyze_image 看图）。
"""
from __future__ import annotations

import csv
import json
from io import StringIO
from pathlib import Path
from typing import Callable

from pydantic import AliasChoices, BaseModel, ConfigDict, Field

from .base import Tool, ToolContext, ToolResult

_PATH_ALIASES = AliasChoices("file_path", "path", "filepath")
_LENIENT = ConfigDict(populate_by_name=True)

# 限制：避免一次性把巨大文件全塞 LLM 上下文。LLM 必须用 offset/limit 翻页。
DEFAULT_LINE_LIMIT = 2000
MAX_LINE_LIMIT = 10000
MAX_BYTES = 50 * 1024 * 1024   # 单文件 50MB 上限，超出直接拒


class ReadFileArgs(BaseModel):
    model_config = _LENIENT
    file_path: str = Field(
        validation_alias=_PATH_ALIASES,
        description="要读取的文件绝对路径。按扩展名自动分派解析器。",
    )
    offset: int = Field(default=0, ge=0, description="起始行号 (0-based)")
    limit: int = Field(
        default=DEFAULT_LINE_LIMIT,
        description=f"最多读取行数，-1 表示默认上限 {DEFAULT_LINE_LIMIT}",
    )
    sheet_name: str | None = Field(
        default=None,
        description="Excel 文件：指定 sheet 名。不传则读第一个 sheet。",
    )
    sheet_index: int | None = Field(
        default=None,
        description="Excel 文件：指定 sheet 下标（0-based）。优先级低于 sheet_name。",
    )
    read_all_sheets: bool = Field(
        default=False,
        description="Excel 文件：读所有 sheet 并合并输出（带分隔符）。",
    )
    complex_mode: bool = Field(
        default=False,
        description=(
            "Excel 复杂表格两阶段读取模式。当默认输出表头混乱/列错位/读不通顺时启用。"
            "第一轮返回原始数据 + [ANNOTATION_REQUIRED] 提示；"
            "按提示生成 segments JSON 后，第二轮通过 annotation 回传得到最终结果。"
        ),
    )
    annotation: str | None = Field(
        default=None,
        description=(
            "complex_mode 第二轮：回传上一轮 [ANNOTATION_REQUIRED] 提示中要求的 "
            "segments JSON 字符串，工具将按此重新编排输出。"
        ),
    )


# ---------------- 各格式解析器 ----------------


def _parse_pdf(p: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(p))
    parts: list[str] = []
    for i, page in enumerate(reader.pages):
        try:
            text = page.extract_text() or ""
        except Exception as e:  # noqa: BLE001
            text = f"[第 {i + 1} 页解析失败: {e}]"
        parts.append(f"## Page {i + 1}\n\n{text.strip()}")
    return "\n\n".join(parts) if parts else "[空 PDF]"


def _parse_docx(p: Path) -> str:
    from docx import Document
    doc = Document(str(p))
    lines: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style:
            lines.append(f"# {text}")
        elif "heading 2" in style:
            lines.append(f"## {text}")
        elif "heading 3" in style:
            lines.append(f"### {text}")
        else:
            lines.append(text)
    # 表格也吐出（每行 | 分隔）
    for ti, table in enumerate(doc.tables):
        lines.append(f"\n### Table {ti + 1}\n")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(lines) if lines else "[空 DOCX]"


def _parse_pptx(p: Path) -> str:
    from pptx import Presentation  # type: ignore[import-untyped]
    prs = Presentation(str(p))
    parts: list[str] = []
    for si, slide in enumerate(prs.slides, start=1):
        slide_parts = [f"## Slide {si}"]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    slide_parts.append(text)
        parts.append("\n\n".join(slide_parts))
    return "\n\n---\n\n".join(parts) if parts else "[空 PPTX]"


def _parse_xlsx_sheet(ws, sheet_name: str) -> str:
    rows: list[list[str]] = []
    for row in ws.iter_rows(values_only=True):
        rows.append(["" if v is None else str(v) for v in row])
    if not rows:
        return f"## Sheet: {sheet_name}\n\n[空]"
    # 用 Markdown 表格输出
    out = [f"## Sheet: {sheet_name}\n"]
    out.append("| " + " | ".join(rows[0]) + " |")
    out.append("| " + " | ".join("---" for _ in rows[0]) + " |")
    for r in rows[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


def _parse_xlsx_complex_raw(ws, sheet_name: str) -> str:
    """complex_mode 第一轮：输出带行号的原始数据，附 ANNOTATION_REQUIRED 指引。"""
    rows: list[list[str]] = []
    for row in ws.iter_rows(values_only=True):
        rows.append(["" if v is None else str(v) for v in row])
    if not rows:
        return f"## Sheet: {sheet_name}\n\n[空]\n\n[ANNOTATION_REQUIRED: 该 sheet 为空，无需 annotation]"

    raw_lines = [f"## Sheet: {sheet_name} — 原始数据（带行号）\n"]
    for i, row in enumerate(rows):
        raw_lines.append(f"Row {i}: {row}")

    annotation_guide = (
        "\n\n[ANNOTATION_REQUIRED]\n"
        "上方为原始行数据。请识别表格结构后，回传如下格式的 segments JSON：\n"
        "```json\n"
        '[\n'
        '  {"header_rows": [0, 1], "data_start_row": 2, "label": "主表"},\n'
        '  {"header_rows": [10], "data_start_row": 11, "label": "子表（可选）"}\n'
        "]\n"
        "```\n"
        "- header_rows: 表头所在行号列表（支持多级表头）\n"
        "- data_start_row: 数据起始行号\n"
        "- label: 该子表的描述（可选）\n"
        "生成完毕后，以 annotation=<上方 JSON 字符串> 重新调用 read_file 得到最终结果。"
    )
    return "\n".join(raw_lines) + annotation_guide


def _parse_xlsx_with_annotation(ws, sheet_name: str, annotation: str) -> str:
    """complex_mode 第二轮：按 annotation segments 重新编排输出。"""
    import json as _json
    rows: list[list[str]] = []
    for row in ws.iter_rows(values_only=True):
        rows.append(["" if v is None else str(v) for v in row])

    try:
        segments = _json.loads(annotation)
    except _json.JSONDecodeError as e:
        return f"[annotation JSON 解析失败: {e}]\n\n回退到标准输出：\n\n" + _parse_xlsx_sheet(ws, sheet_name)

    if not isinstance(segments, list):
        segments = [segments]

    parts: list[str] = []
    for seg in segments:
        header_rows: list[int] = seg.get("header_rows", [0])
        data_start: int = seg.get("data_start_row", header_rows[-1] + 1 if header_rows else 1)
        label: str = seg.get("label", sheet_name)

        if not header_rows or header_rows[0] >= len(rows):
            continue

        # 合并多级表头（水平拼接）
        if len(header_rows) == 1:
            headers = rows[header_rows[0]]
        else:
            merged: list[str] = []
            for col in range(max(len(rows[r]) for r in header_rows)):
                cell_parts = []
                for r in header_rows:
                    v = rows[r][col] if col < len(rows[r]) else ""
                    if v and v not in cell_parts:
                        cell_parts.append(v)
                merged.append("/".join(cell_parts))
            headers = merged

        out = [f"## {label}\n"]
        out.append("| " + " | ".join(str(h) for h in headers) + " |")
        out.append("| " + " | ".join("---" for _ in headers) + " |")
        for ri in range(data_start, len(rows)):
            row = rows[ri]
            padded = list(row) + [""] * max(0, len(headers) - len(row))
            out.append("| " + " | ".join(str(v) for v in padded[:len(headers)]) + " |")
        parts.append("\n".join(out))

    return "\n\n---\n\n".join(parts) if parts else _parse_xlsx_sheet(ws, sheet_name)


def _parse_xlsx(p: Path, *, sheet_name: str | None, sheet_index: int | None,
                  read_all: bool, complex_mode: bool = False,
                  annotation: str | None = None) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(p), read_only=False, data_only=True)
    sheet_names = wb.sheetnames
    if read_all:
        return "\n\n---\n\n".join(
            _parse_xlsx_sheet(wb[name], name) for name in sheet_names
        )
    target: str
    if sheet_name and sheet_name in sheet_names:
        target = sheet_name
    elif sheet_index is not None and 0 <= sheet_index < len(sheet_names):
        target = sheet_names[sheet_index]
    else:
        target = sheet_names[0]

    if complex_mode:
        if annotation:
            return _parse_xlsx_with_annotation(wb[target], target, annotation)
        return _parse_xlsx_complex_raw(wb[target], target)
    return _parse_xlsx_sheet(wb[target], target)


def _parse_csv(p: Path) -> str:
    """CSV 转 Markdown 表格。"""
    with p.open("r", encoding="utf-8", errors="replace", newline="") as f:
        reader = csv.reader(f)
        rows = [row for row in reader]
    if not rows:
        return "[空 CSV]"
    buf = StringIO()
    buf.write("| " + " | ".join(rows[0]) + " |\n")
    buf.write("| " + " | ".join("---" for _ in rows[0]) + " |\n")
    for r in rows[1:]:
        buf.write("| " + " | ".join(r) + " |\n")
    return buf.getvalue().rstrip()


def _parse_plain(p: Path) -> str:
    return p.read_text(encoding="utf-8", errors="replace")


def _parse_json(p: Path) -> str:
    """JSON 文件：解析后用 indent=2 重新输出，便于阅读。失败时退回原文。"""
    raw = p.read_text(encoding="utf-8", errors="replace")
    try:
        parsed = json.loads(raw)
        return json.dumps(parsed, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        return raw


# 扩展名 → 解析器
_PARSERS: dict[str, Callable[..., str]] = {
    ".pdf": lambda p, _a: _parse_pdf(p),
    ".docx": lambda p, _a: _parse_docx(p),
    ".pptx": lambda p, _a: _parse_pptx(p),
    ".xlsx": lambda p, a: _parse_xlsx(
        p, sheet_name=a.sheet_name, sheet_index=a.sheet_index,
        read_all=a.read_all_sheets,
        complex_mode=getattr(a, "complex_mode", False),
        annotation=getattr(a, "annotation", None),
    ),
    ".xlsm": lambda p, a: _parse_xlsx(
        p, sheet_name=a.sheet_name, sheet_index=a.sheet_index,
        read_all=a.read_all_sheets,
        complex_mode=getattr(a, "complex_mode", False),
        annotation=getattr(a, "annotation", None),
    ),
    ".csv": lambda p, _a: _parse_csv(p),
    ".md": lambda p, _a: _parse_plain(p),
    ".markdown": lambda p, _a: _parse_plain(p),
    ".txt": lambda p, _a: _parse_plain(p),
    ".log": lambda p, _a: _parse_plain(p),
    ".json": lambda p, _a: _parse_json(p),
    ".yaml": lambda p, _a: _parse_plain(p),
    ".yml": lambda p, _a: _parse_plain(p),
}


def _apply_pagination(text: str, args: ReadFileArgs) -> tuple[str, int]:
    """按行号切片。返回 (切片后的文本, 总行数)。"""
    lines = text.splitlines()
    total = len(lines)
    limit = args.limit if args.limit > 0 else DEFAULT_LINE_LIMIT
    limit = min(limit, MAX_LINE_LIMIT)
    sliced = lines[args.offset : args.offset + limit]
    return "\n".join(sliced), total


class ReadFileTool(Tool):
    name = "read_file"
    description = (
        "读取复杂格式文档（PDF / DOCX / PPTX / XLSX / CSV / MD / TXT / JSON / YAML），"
        "返回 LLM 友好的 Markdown 文本。支持 offset/limit 分页避免一次塞爆上下文；"
        "Excel 可通过 sheet_name / sheet_index / read_all_sheets 选择 sheet；"
        "遇到复杂 Excel（合并单元格/多级表头）时用 complex_mode=true 两阶段读取。"
    )
    args_model = ReadFileArgs
    risk_level = "low"
    available_to = ("main", "file-agent")

    async def execute(self, args: ReadFileArgs, ctx: ToolContext) -> ToolResult:
        decision = ctx.security.check(tool=self, tool_name=self.name,
                                      args=args.model_dump())
        if decision.action == "block":
            return ToolResult(error=f"risk_blocked: {decision.reason}")
        if decision.action == "confirm":
            return ToolResult(error=f"requires_confirm: {decision.reason}")
        p = Path(args.file_path).expanduser()
        if not p.exists():
            return ToolResult(error=f"文件不存在: {p}")
        if not p.is_file():
            return ToolResult(error=f"不是文件（可能是目录）: {p}")
        size = p.stat().st_size
        if size > MAX_BYTES:
            mb = size / 1024 / 1024
            return ToolResult(error=(
                f"文件过大: {mb:.1f}MB > 50MB 上限。请先用 shell_executor "
                "切片或派 file-agent 处理。"
            ))
        ext = p.suffix.lower()
        parser = _PARSERS.get(ext)
        if parser is None:
            return ToolResult(error=(
                f"不支持的扩展名 '{ext}'。支持: {', '.join(sorted(_PARSERS))}。"
                f"图片请用 analyze_image，其他二进制请考虑 convert_file。"
            ))
        try:
            raw = parser(p, args)
        except Exception as e:  # noqa: BLE001
            return ToolResult(error=f"解析 {ext} 失败: {e}")
        sliced, total = _apply_pagination(raw, args)
        suffix = ""
        if total > args.offset + (args.limit if args.limit > 0 else DEFAULT_LINE_LIMIT):
            suffix = (
                f"\n\n... [已截断，文件共 {total} 行；"
                f"用 offset={args.offset + len(sliced.splitlines())} 继续]"
            )
        return ToolResult(content=sliced + suffix)
